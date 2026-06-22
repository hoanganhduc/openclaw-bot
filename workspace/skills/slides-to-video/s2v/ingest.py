"""Ingest prepared slides (PNG / PDF / PPTX) into an ordered Deck.

All formats converge on one high-resolution PNG per slide plus seed text used to
draft narration. Heavy parsers (Pillow, PyMuPDF, python-pptx, LibreOffice) are
imported lazily so this module loads under base Python.

  PNG  : the frames already exist; order by natural-numeric filename; the
         transcript is seeded by the agent's vision read (no embedded text).
  PDF  : PyMuPDF renders pages at a DPI chosen to hit the target height, and
         extracts page text as a seed.
  PPTX : python-pptx reads on-slide text + speaker notes (best seed); Microsoft
         PowerPoint on Windows or LibreOffice renders the deck to PDF, then the
         PDF path rasterizes the frames.
"""

from __future__ import annotations

import re
import shutil
from pathlib import Path

from . import pptx_render
from .model import Deck, SlideRecord

_NUM_RE = re.compile(r"(\d+)")


def natural_key(name: str) -> list:
    return [int(t) if t.isdigit() else t.lower() for t in _NUM_RE.split(name)]


def detect_format(source: Path) -> str:
    source = Path(source)
    if source.is_dir():
        return "png"
    suffix = source.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix == ".pptx":
        return "pptx"
    if suffix in (".png",):
        return "png"
    raise ValueError(f"unsupported input {source!r}; expected a PNG directory, .pdf, or .pptx")


def _target_dims(resolution: str) -> tuple[int, int]:
    w, h = resolution.lower().split("x", 1)
    return int(w), int(h)


def ingest(source: Path, work_dir: Path, resolution: str = "1920x1080") -> Deck:
    source = Path(source)
    fmt = detect_format(source)
    slides_dir = Path(work_dir) / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)
    if fmt == "png":
        return _ingest_png(source, slides_dir)
    if fmt == "pdf":
        return _ingest_pdf(source, slides_dir, resolution)
    return _ingest_pptx(source, slides_dir, work_dir, resolution)


def _validate_png(path: Path) -> tuple[int, int]:
    from PIL import Image

    Image.MAX_IMAGE_PIXELS = 200_000_000  # decompression-bomb guard
    with Image.open(path) as probe:
        probe.verify()  # structural check
    with Image.open(path) as img:
        img.load()      # force full decode (verify() alone passes some corrupt files)
        return img.width, img.height


def _ingest_png(source: Path, slides_dir: Path) -> Deck:
    files = sorted(source.glob("*.png"), key=lambda p: natural_key(p.name)) if source.is_dir() else [source]
    if not files:
        raise ValueError(f"no PNG files found in {source!r}")
    slides = []
    for i, src in enumerate(files):
        width, height = _validate_png(src)
        dst = slides_dir / f"slide_{i:04d}.png"
        shutil.copyfile(src, dst)
        slides.append(SlideRecord(index=i, image_path=str(dst), source="png",
                                  width=width, height=height,
                                  flags={"needs_vision_seed": True}))
    return Deck(source_format="png", slides=slides)


def _ingest_pdf(source: Path, slides_dir: Path, resolution: str) -> Deck:
    import fitz  # PyMuPDF

    _, target_h = _target_dims(resolution)
    slides = []
    with fitz.open(source) as doc:
        for i, page in enumerate(doc):
            inches_h = page.rect.height / 72.0 or 7.5
            dpi = max(96, int(round(target_h / inches_h)))
            pix = page.get_pixmap(dpi=dpi)
            dst = slides_dir / f"slide_{i:04d}.png"
            pix.save(str(dst))
            text = page.get_text("text").strip()
            slides.append(SlideRecord(index=i, image_path=str(dst), source="pdf",
                                      width=pix.width, height=pix.height,
                                      seed_text=text,
                                      flags={"needs_vision_seed": not bool(text)}))
    return Deck(source_format="pdf", slides=slides)


def _ingest_pptx(source: Path, slides_dir: Path, work_dir: Path, resolution: str) -> Deck:
    from pptx import Presentation

    seeds: list[str] = []
    prs = Presentation(str(source))
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        seeds.append((notes + "\n\n" + "\n".join(parts)).strip())

    pdf_path = pptx_render.pptx_to_pdf(source, Path(work_dir))
    deck = _ingest_pdf(pdf_path, slides_dir, resolution)
    for rec, seed in zip(deck.slides, seeds):  # prefer the richer PPTX text/notes seed
        if seed:
            rec.seed_text = seed
            rec.flags["needs_vision_seed"] = False
        rec.source = "pptx"
    deck.source_format = "pptx"
    return deck
